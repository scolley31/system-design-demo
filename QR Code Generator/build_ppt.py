"""產生 QR Code Generator 系統設計簡報（.pptx）。

從 DESIGN.md 的決策、API、架構與流程整理成投影片。
執行：./.venv/bin/python build_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- palette ----
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
CLIENT = RGBColor(0x64, 0x74, 0x8B)
GATEWAY = RGBColor(0x0E, 0xA5, 0xE9)
APP = RGBColor(0x25, 0x63, 0xEB)
DB = RGBColor(0x05, 0x96, 0x69)
CACHE = RGBColor(0xDC, 0x26, 0x26)
CDN = RGBColor(0x7C, 0x3A, 0xED)
ANALYTICS = RGBColor(0xD9, 0x77, 0x06)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
CJK = "PingFang TC"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _font(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = CJK
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {"typeface": CJK})
    rPr.append(ea)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold, *rest) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(rest[0] if rest else 4)
        r = p.add_run()
        r.text = text
        _font(r, size, color, bold)
    return tb


def box(slide, x, y, w, h, text, fill, fg=WHITE, size=12, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = fill
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ln
        _font(r, size, fg, bold if i == 0 else False)
    return sp


def arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.75):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def header(slide, kicker, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    textbox(slide, Inches(0.55), Inches(0.35), Inches(12), Inches(1.1),
            [(kicker, 13, ACCENT, True, 2), (title, 30, INK, True)])


def bullets(slide, x, y, w, h, items, size=15, gap=7):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = lvl
        r = p.add_run()
        prefix = "•  " if lvl == 0 else "–  "
        r.text = prefix + it
        _font(r, size if lvl == 0 else size - 1, INK if lvl == 0 else MUTED, False)
    return tb


# ============ Slide 1: Title ============
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background(); bg.shadow.inherit = False
box(s, Inches(0.9), Inches(2.0), Inches(1.5), Inches(1.5), "▣", ACCENT, WHITE, 54)
textbox(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(2.5), [
    ("QR Code Generator", 44, WHITE, True, 6),
    ("System Design — 架構、API 與設計決策", 20, RGBColor(0xCB, 0xD5, 0xE1), False, 10),
    ("Dynamic QR · FastAPI 原型 · 對齊 PDF 設計題 + repo 練習", 14, MUTED, False),
])

# ============ Slide 2: Requirements ============
s = prs.slides.add_slide(BLANK)
header(s, "01 · REQUIREMENTS", "需求與規模")
cols = [
    ("功能性 (FR)", APP, [
        "提交 URL → 回 token + QR 圖",
        "掃描短碼 → 302 導向原始 URL",
        "可改目標 URL / 軟刪 / 設過期",
        "掃描分析（總數、每日）",
        "URL 驗證：正規化 + 惡意阻擋",
    ]),
    ("非功能性 (NFR)", DB, [
        "24/7 高可用 (HA)",
        "Redirect 延遲 < 100ms",
        "10 億 QR、1 億用戶",
        "讀多寫少 ≈ 100:1",
    ]),
    ("容量估算", ANALYTICS, [
        "1B × ~200B ≈ 200GB",
        "≈ 5,800 redirect / 秒",
        "1 億用戶 × 5 次/日",
        "圖片 + redirect 為主流量",
    ]),
]
cw = Inches(3.9); gap = Inches(0.25); x0 = Inches(0.6); y0 = Inches(1.8)
for i, (title, color, items) in enumerate(cols):
    x = x0 + i * (cw + gap)
    box(s, x, y0, cw, Inches(0.6), title, color, WHITE, 16)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0 + Inches(0.7), cw, Inches(4.4))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = LIGHT; card.shadow.inherit = False
    bullets(s, x + Inches(0.2), y0 + Inches(0.9), cw - Inches(0.4), Inches(4.0), items, size=14, gap=10)

# ============ Slide 3: API ============
s = prs.slides.add_slide(BLANK)
header(s, "02 · API DESIGN", "API 設計")
rows = [
    ("方法", "路徑", "說明"),
    ("GET", "/", "管理頁前端"),
    ("POST", "/api/v1/qr/create", "建立 → {token, short_url, qr_code_url, original_url}"),
    ("GET", "/r/{token}", "掃描入口：302 / 404(不存在) / 410(刪除·過期)"),
    ("GET", "/api/v1/qr", "列出所有 QR（含掃描數，管理頁）"),
    ("GET", "/api/v1/qr/{token}", "取得 metadata"),
    ("PATCH", "/api/v1/qr/{token}", "改 url / expires_at（部分更新）"),
    ("DELETE", "/api/v1/qr/{token}", "軟刪除"),
    ("GET", "/api/v1/qr/{token}/image", "QR PNG（dimension·color·border）"),
    ("GET", "/api/v1/qr/{token}/analytics", "掃描統計"),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.0)).table
tbl.columns[0].width = Inches(1.4); tbl.columns[1].width = Inches(4.0); tbl.columns[2].width = Inches(6.7)
for r in range(len(rows)):
    for c in range(3):
        cell = tbl.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = rows[r][c]
        _font(run, 12 if r else 13, WHITE if r == 0 else (ACCENT if c == 0 and r else INK), r == 0 or (c == 0))

# ============ Slide 4: Data model ============
s = prs.slides.add_slide(BLANK)
header(s, "03 · DATA MODEL", "資料模型（時間一律存 UTC）")
box(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(0.55), "url_mappings", APP, WHITE, 16)
c1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.45), Inches(5.9), Inches(4.0))
c1.fill.solid(); c1.fill.fore_color.rgb = LIGHT; c1.line.color.rgb = LIGHT; c1.shadow.inherit = False
bullets(s, Inches(0.85), Inches(2.65), Inches(5.4), Inches(3.6), [
    "id  (PK)",
    "token  String(8)  UNIQUE · index",
    "original_url  Text",
    "created_at / updated_at",
    "expires_at  (nullable)",
    "is_deleted  Boolean  ← 軟刪除",
], size=14, gap=10)
box(s, Inches(6.85), Inches(1.8), Inches(5.9), Inches(0.55), "scan_events  (分析明細)", ANALYTICS, WHITE, 16)
c2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(2.45), Inches(5.9), Inches(4.0))
c2.fill.solid(); c2.fill.fore_color.rgb = LIGHT; c2.line.color.rgb = LIGHT; c2.shadow.inherit = False
bullets(s, Inches(7.1), Inches(2.65), Inches(5.4), Inches(3.6), [
    "id (PK) · token · scanned_at",
    "user_agent · ip_address",
    "index (token, scanned_at)",
    "正式版 → 獨立分析庫 + 預聚合",
    "daily_counts(token, date, count)",
], size=14, gap=10)

# ============ Slide: DB Schema ============
def codebox(slide, x, y, w, h, code, size=11):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = LIGHT; sp.line.color.rgb = LIGHT; sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_right = Pt(6); tf.margin_top = Pt(8)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = line if line else " "
        r.font.size = Pt(size); r.font.name = "Menlo"
        r.font.color.rgb = MUTED if line.strip().startswith("--") else INK
    return sp

s = prs.slides.add_slide(BLANK)
header(s, "03b · DB SCHEMA", "資料庫 Schema (DDL) · schema.sql")
ddl1 = """CREATE TABLE url_mappings (
  id            INTEGER  PRIMARY KEY,
  token         VARCHAR(8)   NOT NULL,
  original_url  TEXT     NOT NULL,
  created_at    DATETIME NOT NULL,
  updated_at    DATETIME NOT NULL,
  expires_at    DATETIME     NULL,
  is_deleted    BOOLEAN  NOT NULL DEFAULT 0
);
CREATE UNIQUE INDEX ix_url_mappings_token
  ON url_mappings (token);"""
ddl2 = """CREATE TABLE scan_events (
  id          INTEGER  PRIMARY KEY,
  token       VARCHAR(8)   NOT NULL,
  scanned_at  DATETIME NOT NULL,
  user_agent  VARCHAR(500) NULL,
  ip_address  VARCHAR(45)  NULL
);
CREATE INDEX idx_token_scanned
  ON scan_events (token, scanned_at);"""
codebox(s, Inches(0.6), Inches(1.75), Inches(6.05), Inches(3.5), ddl1)
codebox(s, Inches(6.9), Inches(1.75), Inches(5.85), Inches(3.5), ddl2)
bullets(s, Inches(0.6), Inches(5.45), Inches(12.3), Inches(1.7), [
    "token：8 碼 Base62；UNIQUE 索引兼碰撞安全網（第 5 題）+ redirect O(log n) 查找（第 2 題）",
    "is_deleted 軟刪除（第 12 題）· expires_at 惰性過期（第 13 題）· 時間一律存 UTC",
    "scan_events：分析明細 + 複合索引；正式版改獨立分析庫 + 預聚合（第 14 題）",
], size=13, gap=6)

# ============ Slide 5: Architecture ============
s = prs.slides.add_slide(BLANK)
header(s, "04 · HIGH-LEVEL ARCHITECTURE", "高階架構")
y = Inches(2.3); bh = Inches(0.95)
b_client = box(s, Inches(0.6), y, Inches(1.8), bh, "Client\n掃描者/瀏覽器", CLIENT, WHITE, 13)
b_gw = box(s, Inches(2.9), y, Inches(1.8), bh, "API Gateway", GATEWAY, WHITE, 13)
b_app = box(s, Inches(5.2), y, Inches(2.4), bh, "QR Service\n(stateless, 水平擴展)", APP, WHITE, 13)
arrow(s, Inches(2.4), y + Inches(0.47), Inches(2.9), y + Inches(0.47))
arrow(s, Inches(4.7), y + Inches(0.47), Inches(5.2), y + Inches(0.47))
# right-side stores
rx = Inches(8.6); rw = Inches(4.1)
b_cache = box(s, rx, Inches(1.5), rw, Inches(0.8), "Redis  快取\n(redirect 熱路徑, TTL)", CACHE, WHITE, 12)
b_db = box(s, rx, Inches(2.5), rw, Inches(0.9), "Primary DB → Read Replicas\n(寫走主 · 讀走副 · failover)", DB, WHITE, 12)
b_cdn = box(s, rx, Inches(3.55), rw, Inches(0.8), "Object Store + CDN\n(QR 圖片靜態內容)", CDN, WHITE, 12)
b_an = box(s, rx, Inches(4.45), rw, Inches(0.8), "佇列 → 分析庫\n(非同步 scan · 預聚合)", ANALYTICS, WHITE, 12)
ax = Inches(7.6); ay = y + Inches(0.47)
for tb_, ty in [(b_cache, Inches(1.9)), (b_db, Inches(2.95)), (b_cdn, Inches(3.95)), (b_an, Inches(4.85))]:
    arrow(s, ax, ay, rx, ty)
textbox(s, Inches(0.6), Inches(6.0), Inches(12), Inches(1.0), [
    ("讀多寫少：redirect 先打 Redis，miss 才到 DB（read replica）；QR 圖片走 CDN；scan 非同步進分析庫。", 13, MUTED, False)])

# ============ Slide: Decision map ============
s = prs.slides.add_slide(BLANK)
header(s, "決策地圖", "三大關鍵架構決策（環環相扣）")
cards = [
    ("① 動態 vs 靜態", APP, "選：動態 QR", [
        "QR 編 short URL → server redirect",
        "可修改目標 + 可追蹤掃描",
        "代價：server = SPOF",
        "→ 逼出 cache / CDN / monitoring",
    ]),
    ("② Token 生成策略", DB, "SHA-256+nonce+Base62", [
        "隨機（碰撞靠運氣）",
        "→ 純 hash（同 URL 同 token）",
        "→ +nonce（確定性重試 + DB 兜底）",
        "高流量可換 Pre-gen Pool",
    ]),
    ("③ 301 vs 302", CACHE, "選：302（暫時）", [
        "301 永久快取 → 跳過 server",
        "→ 分析流失、改不了目標",
        "302 每次回源 → 可改/刪/分析",
        "代價：latency → cache + CDN",
    ]),
]
cw = Inches(3.9); gp = Inches(0.3); x0 = Inches(0.5); y = Inches(2.1)
for i, (title, color, choice, items) in enumerate(cards):
    x = x0 + i * (cw + gp)
    box(s, x, y, cw, Inches(0.6), title, color, WHITE, 15)
    box(s, x, y + Inches(0.68), cw, Inches(0.55), choice, GREEN, WHITE, 13)
    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(1.32), cw, Inches(2.7))
    body.fill.solid(); body.fill.fore_color.rgb = LIGHT; body.line.color.rgb = LIGHT; body.shadow.inherit = False
    bullets(s, x + Inches(0.2), y + Inches(1.48), cw - Inches(0.4), Inches(2.4), items, size=12, gap=8)
    if i < 2:
        arrow(s, x + cw, y + Inches(0.3), x + cw + gp, y + Inches(0.3))
textbox(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.6), [
    ("①決定要 server、②決定短碼怎麼來、③決定 redirect 怎麼回 —— 三者串成主架構，其下才是 16 項實作決策。", 13, MUTED, False)])

# ============ Slide 6: Create flow ============
s = prs.slides.add_slide(BLANK)
header(s, "05 · CREATE FLOW", "建立流程")
steps = [
    ("① 驗證", "保守正規化\n+ 惡意/SSRF 檢查", APP),
    ("② 生 token", "SHA-256+nonce\n→Base62→8 碼", GATEWAY),
    ("③ 寫 DB", "直接插入\nUNIQUE 重試 ≤10", DB),
    ("④ 暖 cache", "token→URL\n寫進 Redis", CACHE),
    ("⑤ 生圖", "預生成 QR\n→ object store", CDN),
    ("⑥ 回應", "token / short_url\n/ qr_code_url", ANALYTICS),
]
n = len(steps); bw = Inches(1.85); gap = Inches(0.18); x0 = Inches(0.55); y = Inches(2.7)
for i, (t, d, color) in enumerate(steps):
    x = x0 + i * (bw + gap)
    box(s, x, y, bw, Inches(1.5), t + "\n" + d, color, WHITE, 13)
    if i < n - 1:
        arrow(s, x + bw, y + Inches(0.75), x + bw + gap, y + Inches(0.75))
textbox(s, Inches(0.55), Inches(4.7), Inches(12), Inches(1.2), [
    ("關鍵：第 ③ 步「直接插入 + UNIQUE 例外重試」取代 repo 的「先查再插」，消除高並發 race。", 14, INK, True, 6),
    ("token 碰撞在大規模下必然發生，靠 DB UNIQUE + 重試兜底（nonce 只防『同 URL 同 token』，不防碰撞）。", 13, MUTED, False)])

# ============ Slide 7: Redirect flow ============
s = prs.slides.add_slide(BLANK)
header(s, "06 · REDIRECT FLOW", "掃描 / Redirect 流程（< 100ms 關鍵路徑）")
box(s, Inches(0.6), Inches(1.7), Inches(2.3), Inches(0.8), "GET /r/{token}\n掃描者瀏覽器", CLIENT, WHITE, 12)
b_cache = box(s, Inches(3.3), Inches(1.7), Inches(2.6), Inches(0.8), "① 查 Cache (Redis)", CACHE, WHITE, 13)
arrow(s, Inches(2.9), Inches(2.1), Inches(3.3), Inches(2.1))
# hit path
box(s, Inches(3.3), Inches(2.9), Inches(2.6), Inches(0.7), "命中 + 未過期", GREEN, WHITE, 12)
# miss path
b_db = box(s, Inches(6.4), Inches(1.7), Inches(2.7), Inches(0.8), "② miss → 查 DB\n(read replica)", DB, WHITE, 12)
arrow(s, Inches(5.9), Inches(2.1), Inches(6.4), Inches(2.1))
outs = [
    ("找不到", "404", RED, Inches(9.5), Inches(1.35)),
    ("is_deleted", "410", RED, Inches(9.5), Inches(2.25)),
    ("已過期", "410", RED, Inches(9.5), Inches(3.15)),
    ("正常", "③ 回填 cache", GREEN, Inches(9.5), Inches(4.05)),
]
for label, code, color, x, yy in outs:
    box(s, x, yy, Inches(3.3), Inches(0.7), f"{label}  →  {code}", color, WHITE, 12)
    arrow(s, Inches(9.1), Inches(2.1), x, yy + Inches(0.35), color=MUTED, width=1.2)
# scan + 302
box(s, Inches(3.3), Inches(4.9), Inches(5.8), Inches(0.8),
    "④ 排程非同步記 scan (BackgroundTasks) → 不擋跳轉", ANALYTICS, WHITE, 12)
box(s, Inches(3.3), Inches(5.9), Inches(5.8), Inches(0.8),
    "⑤ 回 302  Location: 原始 URL  → 瀏覽器自動跟隨", APP, WHITE, 13)
arrow(s, Inches(4.6), Inches(3.6), Inches(4.6), Inches(4.9), color=GREEN)
arrow(s, Inches(11.1), Inches(4.75), Inches(9.1), Inches(5.3), color=GREEN, width=1.2)

# ============ Slide: Latency / two-layer lookup ============
s = prs.slides.add_slide(BLANK)
header(s, "06b · LATENCY", "redirect < 100ms — 兩層查找 + 索引")
box(s, Inches(0.6), Inches(2.05), Inches(2.4), Inches(0.9), "掃描\nGET /r/{token}", CLIENT, WHITE, 13)
box(s, Inches(3.5), Inches(2.05), Inches(3.35), Inches(0.9), "① Cache (Redis)\nO(1) · <1ms", CACHE, WHITE, 13)
arrow(s, Inches(3.0), Inches(2.5), Inches(3.5), Inches(2.5))
box(s, Inches(3.5), Inches(3.25), Inches(3.35), Inches(0.8), "命中 → 302 ✅ 不碰 DB", GREEN, WHITE, 12)
arrow(s, Inches(5.17), Inches(2.95), Inches(5.17), Inches(3.25), color=GREEN, width=1.2)
box(s, Inches(7.35), Inches(2.05), Inches(5.4), Inches(0.9), "② miss → SELECT WHERE token=?\n走 ix_url_mappings_token · O(log n)", DB, WHITE, 12)
arrow(s, Inches(6.85), Inches(2.5), Inches(7.35), Inches(2.5))
box(s, Inches(7.35), Inches(3.25), Inches(5.4), Inches(0.8), "回填 cache → 302", APP, WHITE, 12)
arrow(s, Inches(10.05), Inches(2.95), Inches(10.05), Inches(3.25), color=MUTED, width=1.2)
bullets(s, Inches(0.6), Inches(4.5), Inches(12.3), Inches(2.6), [
    "索引對比（查 1 個 token、10 億筆）：無索引 = 全表掃描 O(n)，數秒；有索引 = B-tree O(log n)，次毫秒",
    "scan 記錄非同步（BackgroundTasks），不算進這條延遲",
    "三道防線：cache 擋多數讀（第 7 題）+ token 索引讓 miss 也快（第 2 題）+ scan 非同步不擋路（第 8 題）",
    "正式版再加 read replica 分攤 ~5,800 QPS（第 16 題）",
], size=14, gap=8)

# ============ Slide 8: QR image flow ============
s = prs.slides.add_slide(BLANK)
header(s, "07 · QR IMAGE", "取 QR 圖片流程  ·  GET /api/v1/qr/{token}/image")
steps = [
    ("① 驗證 token", "找不到/已軟刪\n→ 404", APP),
    ("② 組短碼", 'short_url =\nBASE_URL+"/r/{token}"', GATEWAY),
    ("③④ 算矩陣", "編碼→糾錯(RS)\n→2D 模組矩陣", DB),
    ("⑤ 算像素", "每模組→方塊\nborder=quiet zone", CDN),
    ("⑥⑦ PNG", "BytesIO\n選填 dimension 縮放", CACHE),
    ("⑧ 串流", "StreamingResponse\nimage/png", ANALYTICS),
]
n = len(steps); bw = Inches(1.9); gap = Inches(0.14); x0 = Inches(0.5); y = Inches(2.4)
for i, (t, d, color) in enumerate(steps):
    x = x0 + i * (bw + gap)
    box(s, x, y, bw, Inches(1.5), t + "\n" + d, color, WHITE, 12)
    if i < n - 1:
        arrow(s, x + bw, y + Inches(0.75), x + bw + gap, y + Inches(0.75))
textbox(s, Inches(0.55), Inches(4.4), Inches(12.2), Inches(2.4), [
    ("⚠ 編進 QR 的是『不變的短碼 /r/{token}』，不是原始 URL —— dynamic QR 關鍵：改目標 URL 不需重生圖。", 14, INK, True, 8),
    ("QR 原理：字串 → byte 編碼 → Reed-Solomon 糾錯（破損仍可掃）→ 排入矩陣（finder/timing/alignment + 遮罩）。", 13, MUTED, False, 6),
    ("圖片是靜態內容 → 原型即時生成；正式版預生成存 object store + CDN（第 15 題）。", 13, MUTED, False, 6),
    ("樣式參數：dimension（像素邊長）· color（hex 黑塊色）· border（quiet zone 模組數）。", 13, MUTED, False)])

# ============ Slide 07b: QR image storage trade-off ============
s = prs.slides.add_slide(BLANK)
header(s, "07b · QR IMAGE", "圖片存哪 — BLOB-in-DB vs Object Store（第 15 題）")
textbox(s, Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.7), [
    ("關鍵性質:QR 編的是不變短碼 /r/{token} → 圖是靜態內容,一次算好即可純讀;經典取捨:二進位大物該不該進 DB?", 13, INK, False, 2)])
srows = [
    ("方案", "優", "劣"),
    ("即時生成（原型）", "零儲存、URL 改圖自動對", "每次吃 CPU,高 QPS → app 變算圖機"),
    ("存 DB（BLOB 欄位）", "一致性簡單、一個 store、有交易", "PNG 撐大表、吃 buffer pool/備份/複寫;DB 難放 CDN 後 → 貴"),
    ("Object Store + CDN（採用）", "水平無限、單價低、邊緣快取最低延遲", "多一個 store、最終一致"),
]
st = s.shapes.add_table(len(srows), 3, Inches(0.6), Inches(2.25), Inches(12.15), Inches(2.4)).table
st.columns[0].width = Inches(3.0); st.columns[1].width = Inches(4.35); st.columns[2].width = Inches(4.8)
for r in range(len(srows)):
    for c in range(3):
        cell = st.cell(r, c)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = srows[r][c]
        _font(run, 12, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
box(s, Inches(0.6), Inches(4.95), Inches(12.15), Inches(1.0),
    "通則:DB 只存指標(object key / URL),bytes 放 object store —— DB 的價值在交易/索引/查詢,不是當檔案伺服器", DB, WHITE, 13)
bullets(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(1.0), [
    "為何不進 DB:圖只會被整包讀、不查內容 → 佔 buffer pool + 備份/複寫頻寬,且 DB 擋不住 CDN 前的流量",
    ("例外:小圖 + 低量 + 想少一個元件 → BLOB-in-DB 偶爾划算。樣式(dimension/color/border)不同視為不同圖,以 (token,樣式) 為 key", 1),
], size=12, gap=6)

# ============ Slide 9: Token decision ============
s = prs.slides.add_slide(BLANK)
header(s, "08 · DEEP DIVE", "決策 1 — Token 生成與碰撞")
bullets(s, Inches(0.7), Inches(1.9), Inches(11.8), Inches(5), [
    "演算法：SHA-256(url + nonce) → Base62 → 取前 8 碼",
    "nonce = 時間戳_重試次數，打破 SHA-256 的確定性",
    ("作用：讓同一個 URL 每次也產生不同 token（第 4 題：不去重）", 1),
    ("⚠ nonce 只解決『確定性』，不解決『碰撞』", 1),
    "碰撞數學：截斷成 8 碼後空間 = 62^8 ≈ 2.18×10^14",
    ("10 億規模下碰撞『必然』發生（預期 ~2,290 次）", 1),
    ("7 碼 → 8 碼：空間 ×62、預期碰撞降到 1/62", 1),
    "安全網：DB token UNIQUE + 直接插入例外重試（≤10）",
    ("每筆插入碰撞率 ~0.0005%，重試幾乎都第 2 次成功", 1),
], size=15, gap=8)

# ============ Slide: Token internals (Hash -> Encode) ============
s = prs.slides.add_slide(BLANK)
header(s, "08b · TOKEN INTERNALS", "Token 實作 — Hash → Encode")
stages = [
    ("url + nonce", "字串", CLIENT),
    (".encode()", "→ bytes (UTF-8)", GATEWAY),
    ("sha256().digest()", "→ 32 bytes 指紋", APP),
    ("int.from_bytes", "→ 256-bit 大整數", DB),
    ("base62 ÷62 取餘", "→ ~43 字元", CDN),
    ("[:8]", "→ token", ANALYTICS),
]
n = len(stages); bw = Inches(1.9); gap = Inches(0.13); x0 = Inches(0.5); y = Inches(2.25)
for i, (t, d, color) in enumerate(stages):
    x = x0 + i * (bw + gap)
    box(s, x, y, bw, Inches(1.3), t + "\n" + d, color, WHITE, 12)
    if i < n - 1:
        arrow(s, x + bw, y + Inches(0.65), x + bw + gap, y + Inches(0.65))
bullets(s, Inches(0.6), Inches(3.95), Inches(12.2), Inches(3), [
    "Hash（攪亂）：sha256(...).digest() 壓成固定 32 bytes 指紋 — 確定性·雪崩·單向（.hexdigest() 則回 64 字元 hex）",
    "Encode（變短碼）：大整數不斷 ÷62 記餘數 → 餘數 0–61 對映 0-9a-zA-Z；reversed 因餘數從最低位算起",
    ("選 Base62 而非 Base64：不含 / + = → URL 安全", 1),
    "截斷：62^43 ≈ 2^256，取前 8 碼 → 空間砍到 62^8（碰撞源於此）",
    "⚠ 兩個 encode：(url+nonce).encode() 是文字→bytes；base62_encode() 是 bytes→文字",
    "範例：3842 → 餘60='Y'、餘61='Z' → reversed → \"ZY\"",
], size=13, gap=7)

# ============ Slide: Token strategy comparison ============
s = prs.slides.add_slide(BLANK)
header(s, "08c · TOKEN 策略", "Token 生成策略比較（演進路徑）")
trows = [
    ("策略", "簡單度", "唯一性", "可預測", "適用場景"),
    ("Auto-increment", "高", "高", "高（危險）", "內部系統"),
    ("隨機字串 / UUID 截短", "中", "中", "低", "低流量"),
    ("純 hash(URL)", "中", "低", "低", "✗ 同 URL 同 token"),
    ("SHA-256+nonce+Base62", "中", "高", "低", "通用方案 ← 採用"),
    ("Pre-generated Pool", "低", "高", "低", "高流量 · 零碰撞"),
    ("Snowflake-like", "低", "高", "中", "分散式系統"),
]
HILITE = RGBColor(0xDB, 0xEA, 0xFE)
tt = s.shapes.add_table(len(trows), 5, Inches(0.6), Inches(1.75), Inches(12.15), Inches(3.6)).table
tt.columns[0].width = Inches(3.6)
for ci in (1, 2, 3):
    tt.columns[ci].width = Inches(1.75)
tt.columns[4].width = Inches(3.3)
for r in range(len(trows)):
    for c in range(5):
        cell = tt.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = ACCENT
        elif r == 4:
            cell.fill.fore_color.rgb = HILITE
        else:
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = trows[r][c]
        bold = (r == 0) or (c == 0) or (r == 4)
        _font(run, 12, WHITE if r == 0 else INK, bold)
bullets(s, Inches(0.6), Inches(5.6), Inches(12.3), Inches(1.6), [
    "演進：隨機（碰撞靠運氣）→ 純 hash（同 URL 同 token、無法獨立追蹤）→ SHA-256+nonce（換 nonce 確定性重試 + DB 兜底）",
    "進階：Pre-generated Pool（寫入零碰撞，適合超高寫入）· Snowflake-like（分散式本地生成）",
    "我們選 SHA-256+nonce 為通用方案；寫入流量再上去可換 Pre-generated Pool",
], size=13, gap=6)

# ============ Slide 9: Redirect-speed decisions ============
s = prs.slides.add_slide(BLANK)
header(s, "09 · DEEP DIVE", "決策 2 — 302、快取、非同步")
left = [
    "302 而非 301",
    ("301 被瀏覽器快取 → 跳過我方 server", 1),
    ("→ 無法改/刪/統計；故選 302 每次回源", 1),
    "Cache：Redis（非行程內 dict）",
    ("stateless 多台共享 + 全域 invalidation", 1),
    ("行程內 dict 會讓 PATCH/DELETE 失效漏清", 1),
    ("含 TTL 作最終一致性保險", 1),
]
right = [
    "Scan 非同步寫",
    ("redirect 是 <100ms + ~5,800 QPS 熱路徑", 1),
    ("分析可容忍延遲/少量遺失 → 不擋跳轉", 1),
    ("原型 BackgroundTasks；正式版佇列+worker", 1),
    "token index",
    ("避免全表掃描，redirect 直接命中", 1),
]
bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(5), left, size=14, gap=8)
bullets(s, Inches(6.9), Inches(1.9), Inches(6.0), Inches(5), right, size=14, gap=8)

# ============ Slide 09-cache-1: Cache strategies overview ============
s = prs.slides.add_slide(BLANK)
header(s, "附錄 · CACHE 策略", "四種 Cache 策略 — 讀寫怎麼走")
textbox(s, Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.75), [
    ("心法:cache 本質是 replication。兩個問題決定策略 →（1）誰是 master（DB 還是 cache）（2）複製是同步還是非同步。", 13, INK, False, 2),
    ("策略沒有絕對對錯,看 write/read ratio、persistency、consistency 選最適合的。", 12, MUTED, False)])
crows = [
    ("策略", "讀（read）", "寫（write）", "master / 複製"),
    ("Cache Aside", "先查 cache,miss → 讀 DB → 回填", "直接寫 DB（app 自己管 cache）", "DB / 惰性"),
    ("Read Through", "只跟 cache 要,cache 自己回源 DB", "直接寫 DB", "DB / 惰性"),
    ("Write Through", "直接讀 cache", "同時寫 cache + DB,兩者成功才算完成", "同步"),
    ("Write Back", "直接讀 cache", "先寫 cache,再非同步刷回 DB", "cache / 非同步"),
]
ct = s.shapes.add_table(len(crows), 4, Inches(0.6), Inches(2.55), Inches(12.15), Inches(3.2)).table
ct.columns[0].width = Inches(2.2); ct.columns[1].width = Inches(3.55)
ct.columns[2].width = Inches(4.2); ct.columns[3].width = Inches(2.2)
for r in range(len(crows)):
    for c in range(4):
        cell = ct.cell(r, c)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = crows[r][c]
        _font(run, 12, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
textbox(s, Inches(0.6), Inches(5.85), Inches(12.2), Inches(1.3), [
    ("Cache Aside / Read Through:master 是 DB、寫只進 DB → cache 惰性補;差別只在「誰負責回源」(app vs cache 元件)。", 12, MUTED, False, 2),
    ("使用時機:模型不同/後補/部分控制/自寫降級 → Cache Aside(app 回源,本專案);想藏回源、應用變薄、有原生 read-through 產品(如 DAX) → Read Through(cache 回源)。", 12, INK, False, 2),
    ("Write Through / Write Back:寫會碰 cache;差別在對 DB 是同步(慢但不丟) 還是非同步(快但可能丟)。", 12, MUTED, False)])

# ============ Slide 09-cache-2: pros/cons + our choice ============
s = prs.slides.add_slide(BLANK)
header(s, "附錄 · CACHE 策略", "優缺點・適用・本專案落點")
prows = [
    ("策略", "優點", "缺點", "適用"),
    ("Cache Aside", "資料不丟;最易實作", "首讀 cache miss;可能不一致", "讀多、可靠優先"),
    ("Read Through", "程式碼簡潔", "同 Cache Aside", "讀多"),
    ("Write Through", "無 cache miss(讀恆新)", "寫入延遲增加", "讀寫都要新鮮"),
    ("Write Back", "無 miss;抗寫入重;減 DB 負荷", "未刷回前可能丟資料", "寫入重、可容忍丟失"),
]
pt = s.shapes.add_table(len(prows), 4, Inches(0.6), Inches(1.7), Inches(12.15), Inches(3.0)).table
pt.columns[0].width = Inches(2.2); pt.columns[1].width = Inches(3.7)
pt.columns[2].width = Inches(3.65); pt.columns[3].width = Inches(2.6)
for r in range(len(prows)):
    for c in range(4):
        cell = pt.cell(r, c)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = prows[r][c]
        _font(run, 12, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
box(s, Inches(0.6), Inches(5.05), Inches(12.15), Inches(1.0),
    "本專案 redirect:讀走 Cache Aside(miss 回填)、改/刪時 write-around=invalidate(見下一頁 09b)", DB, WHITE, 13)
textbox(s, Inches(0.6), Inches(6.25), Inches(12.2), Inches(0.9), [
    ("為何是 Cache Aside 家族:read-heavy(302 遠多於寫)+ 連結/分析不能丟 → 不選 Write Back(未刷回可能丟);", 12, MUTED, False, 2),
    ("不選 Write Through(多數 QR 從沒被掃,寫時全塞 cache 白佔記憶體)。TTL 兜底最終一致。 來源:homuchen.com/posts/databse-chache-strategies", 11, MUTED, False)])

# ============ Slide 09b: Cache write strategy ============
s = prs.slides.add_slide(BLANK)
header(s, "09b · DEEP DIVE", "Cache 寫策略 — 讀 cache-aside，寫 write-around")
rows = [
    ("redirect 讀", "miss → 查 DB → 回填", "cache-aside", DB),
    ("create", "寫 DB → 順手 set（暖）", "write-through（小優化）", APP),
    ("update / delete", "寫 DB →（commit 後）delete", "write-around / invalidate", CACHE),
]
ry = Inches(1.95)
for op, act, mode, color in rows:
    box(s, Inches(0.7), ry, Inches(2.5), Inches(0.7), op, color, WHITE, 13)
    box(s, Inches(3.35), ry, Inches(4.5), Inches(0.7), act, WHITE, INK, 12, bold=False)
    box(s, Inches(8.0), ry, Inches(4.6), Inches(0.7), mode, WHITE, color, 12)
    ry += Inches(0.9)
bullets(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(2.3), [
    "為何寫時 invalidate 而非 write-through：",
    ("read-heavy 且非每筆會被讀 → write-through 白佔記憶體、擠掉熱門 token；write-around 只快取真正被讀的工作集", 1),
    ("一致性更單純：update 直接刪 cache、下次讀從 DB 重建；順序 commit → 再 invalidate，避免 stale 回填 race", 1),
    ("TTL 兜底最終一致；scan 分析寫入是非同步、不進此 cache（附錄 J）", 1),
    "何時反選 write-through：寫完馬上一定會讀（read-your-write 高頻，如 session）；redirect 建立與被掃無強時間關聯 → write-around 勝",
], size=13, gap=7)

# ============ Slide 09c: write-through 要不要 block ============
s = prs.slides.add_slide(BLANK)
header(s, "09c · DEEP DIVE", "若選 write-through：cache 寫要不要 block？")
textbox(s, Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.5), [
    ("子取捨:client 的寫入要不要「等 cache 寫完才回」?", 13, INK, False)])
wtrows = [
    ("維度", "阻塞式（cache 寫在 critical path）", "非阻塞式（DB commit 後背景更新）"),
    ("寫延遲", "DB + cache 兩趟", "≈ 只有 DB"),
    ("一致性", "強 read-after-write（下一讀保證新）", "commit→更新之間有 stale 窗口"),
    ("cache 故障", "cache 慢/掛 → 寫跟著慢/失敗", "不影響寫入成功（stale 到 TTL）"),
    ("主要坑", "partial failure：DB 已 commit 但 cache 寫失敗 → 難處理", "並發亂序：更新亂序落地 → 停舊值,需 versioning/CAS"),
]
wt = s.shapes.add_table(len(wtrows), 3, Inches(0.6), Inches(2.05), Inches(12.15), Inches(2.9)).table
wt.columns[0].width = Inches(1.7); wt.columns[1].width = Inches(5.35); wt.columns[2].width = Inches(5.1)
for r in range(len(wtrows)):
    for c in range(3):
        cell = wt.cell(r, c)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = wtrows[r][c]
        _font(run, 12, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
bullets(s, Inches(0.6), Inches(5.15), Inches(12.2), Inches(1.0), [
    "另兩個一起想:① 先寫誰 → DB-first（source of truth）再寫 cache（cache-first 若 DB 失敗留幻影值）,但必有 commit→cache 空窗;② 並發互蓋 → cache 最終值取決於「誰最後寫 cache」非「誰最後 commit」→ 需版本號/CAS",
], size=12, gap=5)
box(s, Inches(0.6), Inches(6.25), Inches(12.15), Inches(0.85),
    "∴ 本專案選 invalidate:delete 冪等、對順序不敏感、無部分值不一致 → 繞開「block 與否 + partial failure + 並發互蓋」三個麻煩", CACHE, WHITE, 13)

# ============ Slide 10: Reliability / data ============
s = prs.slides.add_slide(BLANK)
header(s, "10 · DEEP DIVE", "決策 3 — 可靠性與資料語意")
bullets(s, Inches(0.7), Inches(1.9), Inches(11.8), Inches(5), [
    "軟刪除（is_deleted）：支撐 410 語意、保留分析、避免 token 回收風險",
    "錯誤語意：404（從未存在） vs 410（曾存在已刪/過期）",
    ("掃描路徑刪除回 410；管理 API 已刪回 404（刻意差異）", 1),
    "惰性過期 + cron 兜底：掃描即時檢查過期；cron 清過期/通知未點擊/物理清理",
    "URL 保守正規化：只 lower-case host，保留 path/query 與原 scheme",
    ("修正 repo 的整串小寫（破壞 path）+ 強制 https（導向失效站）", 1),
    "惡意阻擋：黑名單 + SSRF/內網位址（localhost·127.*·私網），正式版接 Safe Browsing",
], size=15, gap=9)

# ============ Slide 11: Scaling ============
s = prs.slides.add_slide(BLANK)
header(s, "11 · SCALING", "擴展策略")
bullets(s, Inches(0.7), Inches(1.9), Inches(11.8), Inches(5), [
    "App server stateless → 增加 instance 水平擴展（cache 外置 Redis）",
    "DB：單 primary + 多 read replica（讀走副）+ primary 掛 failover 升主",
    ("複寫延遲：create 後暖 cache 緩解", 1),
    "分片：以 token hash 為鍵（均勻），列為未來選項（200GB 目前單機足夠）",
    "QR 圖片：object store + CDN（靜態內容，目標 URL 改不需重生圖）",
    "分析：佇列 → 獨立分析庫 + 預聚合每日計數表",
    "清理 cron：清過期 / 通知後刪長期未點擊 / 物理清理超期軟刪 + invalidate cache",
], size=15, gap=9)

# ============ Slide: Read-heavy -> Read Replica ============
s = prs.slides.add_slide(BLANK)
header(s, "11b · READ REPLICA", "Read-heavy → 讀寫分離 + Read Replica")
textbox(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(3.2), [
    ("容量：1B × 200B = 200GB（單機放得下）", 15, INK, False, 8),
    ("讀流量：1 億 users × 5/日 = 5 億 redirects/日", 15, INK, False, 4),
    ("            5 億 ÷ 86,400s ≈ 5,787 redirects/秒", 16, ACCENT, True, 8),
    ("讀寫比：≈ 100:1（讀 ≫ 寫）", 15, INK, False, 12),
    ("→ 瓶頸是「讀吞吐」、不是容量", 15, INK, True, 4),
    ("→ 讀寫分離 + 多個 read replica", 16, GREEN, True),
])
# topology
b_app = box(s, Inches(6.7), Inches(2.5), Inches(1.8), Inches(1.0), "App\n(stateless)", APP, WHITE, 13)
box(s, Inches(9.4), Inches(1.9), Inches(3.3), Inches(0.8), "Primary (Write)", DB, WHITE, 13)
box(s, Inches(9.4), Inches(2.95), Inches(3.3), Inches(0.65), "Read Replica 1", GATEWAY, WHITE, 12)
box(s, Inches(9.4), Inches(3.75), Inches(3.3), Inches(0.65), "Read Replica 2", GATEWAY, WHITE, 12)
box(s, Inches(9.4), Inches(4.55), Inches(3.3), Inches(0.65), "Read Replica N", GATEWAY, WHITE, 12)
arrow(s, Inches(8.5), Inches(2.8), Inches(9.4), Inches(2.3), color=RED, width=1.6)         # 寫
arrow(s, Inches(8.5), Inches(3.15), Inches(9.4), Inches(3.27), color=GREEN, width=1.4)     # 讀
arrow(s, Inches(8.5), Inches(3.3), Inches(9.4), Inches(4.07), color=GREEN, width=1.4)
arrow(s, Inches(8.5), Inches(3.4), Inches(9.4), Inches(4.87), color=GREEN, width=1.4)
textbox(s, Inches(8.35), Inches(2.25), Inches(1.0), Inches(0.3), [("寫", 12, RED, True)])
textbox(s, Inches(8.3), Inches(3.5), Inches(1.0), Inches(0.3), [("讀", 12, GREEN, True)])
arrow(s, Inches(11.05), Inches(2.7), Inches(11.05), Inches(2.95), color=MUTED, width=1.1)  # 複寫
textbox(s, Inches(11.2), Inches(2.62), Inches(1.6), Inches(0.3), [("非同步複寫", 10, MUTED, False)])
bullets(s, Inches(0.6), Inches(5.5), Inches(12.3), Inches(1.7), [
    "寫走 primary、讀走 replica（load-balanced）；讀不夠就加 replica 水平擴讀",
    "配合 Redis cache，真正打到 DB 的讀已先被擋掉一大半",
    "failover：primary 掛 → 升一台 read replica 為新 primary；複寫延遲靠 create 後暖 cache 緩解",
], size=13, gap=6)

# ============ Slide 12: Improvements over repo ============
s = prs.slides.add_slide(BLANK)
header(s, "12 · IMPROVEMENTS", "對參考 repo 的關鍵改良")
items = [
    ("①", "先查再插 → 直接插入 + UNIQUE 例外重試", "消除高並發 race"),
    ("②", "行程內 dict cache → Redis", "stateless 一致性（invalidation 全域生效）"),
    ("③", "同步寫 scan → 非同步", "不擋 redirect 關鍵路徑"),
    ("④", "整串小寫 → 保守正規化", "避免破壞 path 大小寫 / 導向失效站"),
    ("＋", "token 7 碼 → 8 碼", "碰撞預期 ×1/62"),
]
y = Inches(1.85)
for num, change, why in items:
    box(s, Inches(0.7), y, Inches(0.7), Inches(0.8), num, ACCENT, WHITE, 18)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.55), y, Inches(11.1), Inches(0.8))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = LIGHT; card.shadow.inherit = False
    textbox(s, Inches(1.8), y + Inches(0.08), Inches(10.7), Inches(0.7), [
        (change, 15, INK, True, 2), (why, 12, MUTED, False)])
    y += Inches(1.0)

# ============ Slide: NoSQL alternative ============
s = prs.slides.add_slide(BLANK)
header(s, "附錄 · NOSQL", "替代方案 — DynamoDB 風格 key 設計")
box(s, Inches(0.6), Inches(1.65), Inches(6.0), Inches(0.6),
    "直覺方案  PK=user_id · SK=created_at#qr_token", ANALYTICS, WHITE, 12)
ca = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.3), Inches(6.0), Inches(1.9))
ca.fill.solid(); ca.fill.fore_color.rgb = LIGHT; ca.line.color.rgb = LIGHT; ca.shadow.inherit = False
bullets(s, Inches(0.8), Inches(2.45), Inches(5.6), Inches(1.7), [
    "✅ 列出我的 QR：Query PK=user_id 天生最佳",
    "❌ redirect 只有 token、沒 user_id → 定位不到分區",
    ("→ 退化成 Scan O(n)，違反 <100ms（致命）", 1),
], size=13, gap=6)
box(s, Inches(6.85), Inches(1.65), Inches(5.9), Inches(0.6),
    "建議  base PK=qr_token + GSI(user_id, created_at)", GREEN, WHITE, 12)
cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(2.3), Inches(5.9), Inches(1.9))
cb.fill.solid(); cb.fill.fore_color.rgb = LIGHT; cb.line.color.rgb = LIGHT; cb.shadow.inherit = False
bullets(s, Inches(7.05), Inches(2.45), Inches(5.5), Inches(1.7), [
    "✅ redirect：GetItem(PK=qr_token) O(1)、強一致",
    "✅ 列我的 QR：走 GSI(PK=user_id)",
    "qr_token 高亂度 → 均勻分區，避免 hot partition",
    "唯一性：條件寫入 attribute_not_exists",
], size=13, gap=5)
nrows = [
    ("存取模式", "SQL（現況）", "NoSQL PK=user_id", "NoSQL PK=qr_token"),
    ("redirect 查 token", "B-tree O(log n) ✅", "Scan O(n) ❌", "GetItem O(1) ✅"),
    ("列出我的 QR", "WHERE user_id ✅", "天生最佳 ✅", "GSI 查 ✅"),
    ("token 唯一", "UNIQUE ✅", "需額外處理", "條件寫入 ✅"),
]
nt = s.shapes.add_table(len(nrows), 4, Inches(0.6), Inches(4.5), Inches(12.15), Inches(2.1)).table
nt.columns[0].width = Inches(3.0)
for ci in (1, 2, 3):
    nt.columns[ci].width = Inches(3.05)
for r in range(len(nrows)):
    for c in range(4):
        cell = nt.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = nrows[r][c]
        _font(run, 11, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
textbox(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.6),
        [("NoSQL 先列存取模式、再決定 key；redirect（token→URL）是這個系統的第一公民。", 12, MUTED, False)])

# ============ Slide: CDN edge redirect variant ============
s = prs.slides.add_slide(BLANK)
header(s, "附錄 · CDN EDGE", "CDN 邊緣 redirect 變體（可選最佳化）")
textbox(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.6), [
    ("主設計選擇不快取 redirect（回源保即時改/刪 + 分析）。302 與 CDN 無因果——302 仍可設短 TTL 讓 CDN 快取；", 12, MUTED, False, 2),
    ("真正取捨是「快取 redirect vs 即時性/分析」。本變體：熱門 token 的 302 設短 TTL、邊緣直接跳，換更低延遲。", 12, MUTED, False)])
box(s, Inches(0.6), Inches(2.2), Inches(2.2), Inches(0.95), "掃描者\nGET /r/{token}", CLIENT, WHITE, 12)
box(s, Inches(3.3), Inches(2.1), Inches(3.6), Inches(1.15),
    "CDN Edge（近使用者）\nedge cache: token→URL\n短 TTL（例 60s）", CDN, WHITE, 12)
arrow(s, Inches(2.8), Inches(2.65), Inches(3.3), Inches(2.65))
box(s, Inches(3.3), Inches(3.65), Inches(3.6), Inches(0.85),
    "命中 → CDN 直接回 302\n不回源 · 最低延遲", GREEN, WHITE, 12)
arrow(s, Inches(5.1), Inches(3.25), Inches(5.1), Inches(3.65), color=GREEN, width=1.3)
box(s, Inches(7.3), Inches(2.25), Inches(5.4), Inches(0.95),
    "miss → 回源 Service\n→ Cache/DB → 回 302（CDN 記短 TTL）", DB, WHITE, 12)
arrow(s, Inches(6.9), Inches(2.65), Inches(7.3), Inches(2.65))
nrows = [
    ("面向", "主設計（回源 + Redis）", "CDN 邊緣變體"),
    ("延遲", "低（回源 + Redis）", "最低（邊緣直接跳）"),
    ("改 / 刪即時性", "立即生效", "TTL 內延遲"),
    ("掃描分析", "每次都記", "TTL 內漏記"),
]
nt = s.shapes.add_table(len(nrows), 3, Inches(0.6), Inches(4.75), Inches(12.15), Inches(2.0)).table
nt.columns[0].width = Inches(3.0); nt.columns[1].width = Inches(4.55); nt.columns[2].width = Inches(4.6)
for r in range(len(nrows)):
    for c in range(3):
        cell = nt.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = nrows[r][c]
        _font(run, 12, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
textbox(s, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.5), [
    ("預設不採用；列為超高流量、目標穩定的熱門 QR 之可選最佳化（短 TTL 控制誤差）。", 12, MUTED, False)])

# ============ Slide: Prototype -> Production gap ============
s = prs.slides.add_slide(BLANK)
header(s, "附錄 · PROD GAP", "從 Prototype 到 Production 的落差")
textbox(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.5), [
    ("動態 QR 把 server 變成 SPOF → caching + CDN + monitoring 不是加分項，是必要代價。", 13, MUTED, False)])
grows = [
    ("面向", "原型現況", "Production 需要"),
    ("Error handling", "✅ 已實作", "全域處理 + 降級 + request id"),
    ("Rate limiting", "✅ 已實作", "WAF per-IP + API GW 節流"),
    ("Auth & 多租戶", "✅ 已實作", "Cognito + JWT + owner_id 隔離"),
    ("Monitoring / Alerting", "✅ 已實作", "CloudWatch alarms→SNS + Logs + canary"),
    ("Data cleanup", "✅ 已實作", "EventBridge + Lambda 定時清理"),
    ("Caching / CDN", "記憶體 + 即時生圖", "Redis + object store + CDN"),
]
gt = s.shapes.add_table(len(grows), 3, Inches(0.6), Inches(2.15), Inches(12.15), Inches(4.0)).table
gt.columns[0].width = Inches(3.2); gt.columns[1].width = Inches(4.2); gt.columns[2].width = Inches(4.75)
RED_BG = RGBColor(0xFE, 0xE2, 0xE2)
for r in range(len(grows)):
    for c in range(3):
        cell = gt.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = ACCENT
        elif c == 1 and grows[r][1] == "無":
            cell.fill.fore_color.rgb = RED_BG
        else:
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = grows[r][c]
        _font(run, 12, WHITE if r == 0 else INK, r == 0 or c == 0)
textbox(s, Inches(0.6), Inches(6.35), Inches(12.2), Inches(0.5), [
    ("rate limiting · auth/多租戶 · monitoring 是目前完全沒談、production 前必補的三塊。", 12, MUTED, False)])

# ============ Slide: AWS Deployment Architecture ============
s = prs.slides.add_slide(BLANK)
header(s, "部署 · AWS", "AWS 部署架構（Terraform，已上線）")
box(s, Inches(0.5), Inches(2.7), Inches(1.7), Inches(0.85), "Client", CLIENT, WHITE, 13)
box(s, Inches(2.45), Inches(2.7), Inches(2.0), Inches(0.85), "CloudFront\n(CDN + WAF)", CDN, WHITE, 12)
box(s, Inches(2.45), Inches(1.55), Inches(2.0), Inches(0.7), "S3 (QR PNG)\n/qr-img/*", ANALYTICS, WHITE, 11)
box(s, Inches(4.7), Inches(2.7), Inches(2.0), Inches(0.85), "API Gateway\n(JWT auth · 節流)", GATEWAY, WHITE, 12)
box(s, Inches(4.7), Inches(1.55), Inches(2.0), Inches(0.7), "Cognito\n(user pool)", DB, WHITE, 11)
arrow(s, Inches(5.7), Inches(2.7), Inches(5.7), Inches(2.25), color=MUTED, width=1.2)  # APIGW→Cognito(JWT)
box(s, Inches(6.95), Inches(2.7), Inches(1.6), Inches(0.85), "內部 ALB\n/health", APP, WHITE, 12)
box(s, Inches(8.8), Inches(2.7), Inches(2.1), Inches(0.95), "EC2 ASG\nDocker/gunicorn", INK, WHITE, 12)
box(s, Inches(8.8), Inches(3.95), Inches(2.1), Inches(0.7), "RDS PostgreSQL", DB, WHITE, 11)
box(s, Inches(8.8), Inches(4.8), Inches(2.1), Inches(0.7), "ElastiCache Redis", CACHE, WHITE, 11)
arrow(s, Inches(2.2), Inches(3.12), Inches(2.45), Inches(3.12))
arrow(s, Inches(3.45), Inches(2.7), Inches(3.45), Inches(2.25), color=MUTED, width=1.2)  # CF→S3
arrow(s, Inches(4.45), Inches(3.12), Inches(4.7), Inches(3.12))
arrow(s, Inches(6.7), Inches(3.12), Inches(6.95), Inches(3.12))
arrow(s, Inches(8.55), Inches(3.12), Inches(8.8), Inches(3.12))
arrow(s, Inches(9.85), Inches(3.65), Inches(9.85), Inches(3.95), color=MUTED, width=1.2)
arrow(s, Inches(9.85), Inches(4.65), Inches(9.85), Inches(4.8), color=MUTED, width=1.2)
textbox(s, Inches(4.7), Inches(3.62), Inches(2.2), Inches(0.3), [("VPC Link", 10, MUTED, False)])
textbox(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.4), [
    ("WAF(per-IP rate limit)掛在 CloudFront 邊緣;/qr-img/* → S3 長快取;/, /r/*, /api/* → API Gateway,redirect 不快取。", 13, INK, False, 4),
    ("設定/密鑰:SSM Parameter Store(DATABASE_URL/REDIS_URL/S3_BUCKET/CDN_BASE/BASE_URL/IMAGE_URI)·Secrets Manager(DB 密碼)·ECR·NAT Gateway。", 12, MUTED, False)])

# ============ Slide: Rate Limiting — two layers ============
s = prs.slides.add_slide(BLANK)
header(s, "安全 · RATE LIMITING", "兩層限流 — Application Layer vs Infrastructure Layer")
textbox(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(0.55), [
    ("情境:有人用 script 一秒打你 10 萬次,怎麼擋?兩層各有所長,production 通常縱深並用。", 13, INK, False, 2)])
rlcmp = [
    ("維度", "Application Layer", "Infrastructure Layer"),
    ("常見手段", "token bucket / sliding window、per-API-key quota、回 429", "CDN + WAF（擋 L3–L7 DDoS）、LB connection limit、API Gateway throttling"),
    ("跑在哪", "你的 server code（middleware、sidecar）", "CDN edge、WAF、LB、API Gateway"),
    ("看得到的資訊", "業務 context（user ID、訂閱 tier、資源所有權、業務狀態）", "HTTP 通用內容（IP、headers、URL、body bytes），不解業務語意"),
    ("拒絕成本", "高（已花 TCP+TLS+server CPU）", "低（在 edge 丟掉，不耗 origin）"),
    ("粒度", "細（per user、per endpoint、per resource）", "粗（per IP、per region、per connection）"),
]
ct = s.shapes.add_table(len(rlcmp), 3, Inches(0.6), Inches(2.05), Inches(12.15), Inches(3.9)).table
ct.columns[0].width = Inches(2.1); ct.columns[1].width = Inches(5.0); ct.columns[2].width = Inches(5.05)
for r in range(len(rlcmp)):
    for c in range(3):
        cell = ct.cell(r, c)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = rlcmp[r][c]
        _font(run, 12, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
box(s, Inches(0.6), Inches(6.15), Inches(12.15), Inches(1.0),
    "本專案目前只做 Infrastructure Layer（WAF per-IP + API GW 節流）—— 擋匿名洪水最划算;per-user/tier 的業務限流（App Layer）列為後續", GATEWAY, WHITE, 13)

# ============ Slide: Rate Limiting ============
s = prs.slides.add_slide(BLANK)
header(s, "安全 · RATE LIMITING", "本專案實作：WAF per-IP + API Gateway 節流（Infra 層縱深）")
box(s, Inches(0.5), Inches(2.0), Inches(1.8), Inches(0.9), "掃描者 / Script", CLIENT, WHITE, 12)
box(s, Inches(2.65), Inches(2.0), Inches(2.5), Inches(0.9), "WAF (per-IP)\n超量 → 403", RED, WHITE, 12)
box(s, Inches(5.55), Inches(2.0), Inches(1.9), Inches(0.9), "CloudFront", CDN, WHITE, 12)
box(s, Inches(7.85), Inches(2.0), Inches(2.7), Inches(0.9), "API GW 整體節流\n超量 → 429", GATEWAY, WHITE, 12)
box(s, Inches(10.95), Inches(2.0), Inches(1.85), Inches(0.9), "ALB → EC2", APP, WHITE, 11)
arrow(s, Inches(2.3), Inches(2.45), Inches(2.65), Inches(2.45))
arrow(s, Inches(5.15), Inches(2.45), Inches(5.55), Inches(2.45))
arrow(s, Inches(7.45), Inches(2.45), Inches(7.85), Inches(2.45))
arrow(s, Inches(10.55), Inches(2.45), Inches(10.95), Inches(2.45))
rlrows = [
    ("層", "機制", "上限", "超量"),
    ("CloudFront WAF", "per-IP · /api/*", "300 / 5 分 / IP", "403"),
    ("CloudFront WAF", "per-IP · 全域", "2000 / 5 分 / IP", "403"),
    ("API Gateway", "整體 throttle", "1000 req/s · burst 2000", "429"),
]
rt = s.shapes.add_table(len(rlrows), 4, Inches(0.6), Inches(3.3), Inches(12.15), Inches(2.0)).table
rt.columns[0].width = Inches(3.0); rt.columns[1].width = Inches(3.0); rt.columns[2].width = Inches(4.15); rt.columns[3].width = Inches(2.0)
for r in range(len(rlrows)):
    for c in range(4):
        cell = rt.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = rlrows[r][c]
        _font(run, 12, WHITE if r == 0 else INK, r == 0 or c == 0)
bullets(s, Inches(0.6), Inches(5.55), Inches(12.3), Inches(1.6), [
    "為何不能只靠 API GW:HTTP API 內建節流是『整體』非 per-IP、也無 API key → 擋單一 IP 濫用必須靠 WAF",
    "限速值皆 Terraform 變數可調;WAF 阻擋數進 CloudWatch（BlockedRequests）",
    "純 infra 變更,不動 app;已 validate/plan 通過(未 apply)",
], size=13, gap=7)

# ============ Slide: Auth & Isolation ============
s = prs.slides.add_slide(BLANK)
header(s, "安全 · AUTH & 多租戶", "Cognito + API GW JWT authorizer + owner_id 隔離")
box(s, Inches(0.5), Inches(2.0), Inches(2.2), Inches(0.9), "前端\n(管理頁)", CLIENT, WHITE, 12)
box(s, Inches(3.0), Inches(2.0), Inches(2.6), Inches(0.9), "Cognito Hosted UI\n登入(PKCE)→ id token", DB, WHITE, 11)
box(s, Inches(5.95), Inches(2.0), Inches(2.8), Inches(0.9), "API GW JWT authorizer\n未帶/無效 → 401", GATEWAY, WHITE, 11)
box(s, Inches(9.1), Inches(2.0), Inches(3.4), Inches(0.9), "EC2:取 sub=owner_id\n依 owner 過濾/授權", APP, WHITE, 11)
arrow(s, Inches(2.7), Inches(2.45), Inches(3.0), Inches(2.45))
arrow(s, Inches(5.6), Inches(2.45), Inches(5.95), Inches(2.45))
arrow(s, Inches(8.75), Inches(2.45), Inches(9.1), Inches(2.45))
nrows = [
    ("", "需登入 (JWT)", "公開 (無 auth)"),
    ("端點", "create / list / get / patch / delete / analytics", "/r/{token} 掃描 · QR 圖 · /health"),
    ("隔離", "owner_id = Cognito sub;非擁有者一律 404", "redirect/圖不分 owner"),
]
nt = s.shapes.add_table(len(nrows), 3, Inches(0.6), Inches(3.4), Inches(12.15), Inches(1.7)).table
nt.columns[0].width = Inches(1.6); nt.columns[1].width = Inches(6.35); nt.columns[2].width = Inches(4.2)
for r in range(len(nrows)):
    for c in range(3):
        cell = nt.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = nrows[r][c]
        _font(run, 12, WHITE if r == 0 else INK, r == 0 or c == 0)
bullets(s, Inches(0.6), Inches(5.4), Inches(12.3), Inches(1.6), [
    "前端用 id token(帶 aud+email);API GW authorizer 在邊緣擋未授權,app 再取 sub 當 owner_id",
    "env-gated:本機不設 AUTH_ENABLED → dev user『local-dev』、免登入,SQLite 照跑",
    "純 env 開關 + 新 infra(Cognito/authorizer),已 validate/plan 通過(未 apply)",
], size=13, gap=7)

# ============ Slide: Data Cleanup Cron ============
s = prs.slides.add_slide(BLANK)
header(s, "可靠性 · DATA CLEANUP", "定時清理:EventBridge + Lambda（防 DB bloat）")
box(s, Inches(0.6), Inches(2.1), Inches(3.0), Inches(0.95), "EventBridge Scheduler\n每日 03:00 UTC", ANALYTICS, WHITE, 12)
box(s, Inches(4.0), Inches(2.1), Inches(3.0), Inches(0.95), "Lambda (VPC)\npg8000 純 Python", APP, WHITE, 12)
box(s, Inches(7.4), Inches(2.1), Inches(2.6), Inches(0.95), "RDS PostgreSQL\nDELETE", DB, WHITE, 12)
arrow(s, Inches(3.6), Inches(2.57), Inches(4.0), Inches(2.57))
arrow(s, Inches(7.0), Inches(2.57), Inches(7.4), Inches(2.57))
crows = [
    ("清理目標", "條件", "預設保留"),
    ("過期 QR", "expires_at < now - grace", "30 天"),
    ("軟刪除超期", "is_deleted 且 updated_at < now - retention", "30 天"),
    ("連帶", "上述 token 的 scan_events 一併刪", "—"),
]
ct = s.shapes.add_table(len(crows), 3, Inches(0.6), Inches(3.5), Inches(12.15), Inches(1.9)).table
ct.columns[0].width = Inches(2.6); ct.columns[1].width = Inches(7.05); ct.columns[2].width = Inches(2.5)
for r in range(len(crows)):
    for c in range(3):
        cell = ct.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = crows[r][c]
        _font(run, 12, WHITE if r == 0 else INK, r == 0 or c == 0)
bullets(s, Inches(0.6), Inches(5.6), Inches(12.3), Inches(1.4), [
    "單一來源:cleanup.py(SQLAlchemy+text SQL)Lambda 與本機共用;Lambda 用純 Python pg8000 免二進位打包",
    "與 app 實例解耦(不用 app 內排程,免多實例重複跑/leader election);保留天數皆 Terraform 變數可調",
], size=13, gap=8)

# ============ Slide: Monitoring / Alerting ============
s = prs.slides.add_slide(BLANK)
header(s, "可靠性 · MONITORING", "CloudWatch alarms → SNS + Logs + Dashboard + Canary")
box(s, Inches(0.6), Inches(2.0), Inches(3.2), Inches(0.9), "CloudWatch Alarms\nALB/RDS/EC2/API GW/Lambda/CF", APP, WHITE, 11)
box(s, Inches(4.3), Inches(2.0), Inches(2.6), Inches(0.9), "SNS topic\n→ Email 通知", ANALYTICS, WHITE, 12)
box(s, Inches(7.4), Inches(2.0), Inches(2.6), Inches(0.9), "Synthetic Canary\n每 5 分 /health", DB, WHITE, 12)
arrow(s, Inches(3.8), Inches(2.45), Inches(4.3), Inches(2.45))
arrow(s, Inches(7.4), Inches(2.45), Inches(6.9), Inches(2.45), color=MUTED, width=1.2)  # canary→alarms
mrows = [
    ("面向", "內容"),
    ("告警 alarms", "ALB unhealthy/5xx/latency · RDS CPU/storage · ElastiCache · API GW 5xx · Lambda errors · CloudFront 5xx(us-east-1)"),
    ("通知", "CloudWatch alarm → SNS topic → email（alert_email,需點確認信）"),
    ("App logs", "EC2 容器 awslogs driver → CloudWatch Logs /qrcode/app（retention 14d）"),
    ("Dashboard", "qrcode-overview：ALB/RDS/EC2/API GW/Lambda 一頁"),
    ("Canary", "syn-nodejs-puppeteer 每 5 分探測 CloudFront /health → 失敗即 alarm"),
]
mt = s.shapes.add_table(len(mrows), 2, Inches(0.6), Inches(3.3), Inches(12.15), Inches(2.9)).table
mt.columns[0].width = Inches(2.2); mt.columns[1].width = Inches(9.95)
for r in range(len(mrows)):
    for c in range(2):
        cell = mt.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = mrows[r][c]
        _font(run, 12, WHITE if r == 0 else INK, r == 0 or c == 0)
bullets(s, Inches(0.6), Inches(6.4), Inches(12.3), Inches(0.8), [
    "解決『服務掛了沒人知道』;canary 是端到端外部探測(連 CloudFront/整鏈路掛掉都測得到)。純 infra,已 validate/plan 通過(未 apply)。",
], size=12, gap=4)

# ============ Slide: Error Handling ============
s = prs.slides.add_slide(BLANK)
header(s, "可靠性 · ERROR HANDLING", "壞輸入/依賴故障都優雅回應、不洩漏、可追蹤")
ehrows = [
    ("面向", "做法"),
    ("全域處理器", "未捕捉例外 → 500 不洩漏 + 記 traceback;422 驗證轉可讀字串(沿用 {detail})"),
    ("依賴降級", "Redis 故障 → redirect 走 DB(不 500);S3 失敗 → create fallback /image"),
    ("Request ID", "所有回應帶 X-Request-ID + 進日誌/錯誤 body(對應 CloudWatch Logs)"),
    ("輸入健全化", "body > 64KB → 413;malformed JSON → 422;URL 長度 2048"),
]
et = s.shapes.add_table(len(ehrows), 2, Inches(0.6), Inches(2.0), Inches(12.15), Inches(2.7)).table
et.columns[0].width = Inches(2.4); et.columns[1].width = Inches(9.75)
for r in range(len(ehrows)):
    for c in range(2):
        cell = et.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = ehrows[r][c]
        _font(run, 12, WHITE if r == 0 else INK, r == 0 or c == 0)
bullets(s, Inches(0.6), Inches(4.9), Inches(12.3), Inches(2.0), [
    "錯誤格式沿用 {\"detail\"} → 相容前端;500 只回通用訊息 + request_id,traceback 僅進日誌",
    "依賴降級讓「Redis/S3 掛掉」不再連帶整個請求 500;分析寫入失敗也只記 log",
    "本機驗證:422/413/500(含 X-Request-ID)+ REDIS_URL 不可達時 redirect 仍 302。純 app,不動 infra",
], size=13, gap=8)

# ============ Slide: CI/CD Flow ============
s = prs.slides.add_slide(BLANK)
header(s, "部署 · CI/CD", "CI/CD 流程（GitHub Actions + OIDC）")
steps = [
    ("git push", "QR Code Generator/**", CLIENT),
    ("GitHub Actions", "OIDC assume role\n(無長期金鑰)", GATEWAY),
    ("buildx", "--platform\nlinux/arm64", APP),
    ("ECR push", "tag = SHA\n+ latest", DB),
    ("SSM 更新", "/qrcode/\nIMAGE_URI", CDN),
    ("SSM 部署", "→ EC2\ndeploy-app.sh", ANALYTICS),
    ("上線", "ALB /health\n通過", GREEN),
]
n = len(steps); bw = Inches(1.62); gp = Inches(0.13); x0 = Inches(0.5); y = Inches(2.5)
for i, (t, d, color) in enumerate(steps):
    x = x0 + i * (bw + gp)
    box(s, x, y, bw, Inches(1.4), t + "\n" + d, color, WHITE, 11)
    if i < n - 1:
        arrow(s, x + bw, y + Inches(0.7), x + bw + gp, y + Inches(0.7))
bullets(s, Inches(0.6), Inches(4.4), Inches(12.3), Inches(2.4), [
    "EC2 deploy-app.sh:從 SSM 撈設定 → ECR login → docker pull → docker run(換新容器)",
    "全 env-gated:本機無這些環境變數 → 維持 SQLite + 記憶體 cache + 即時生圖,不受影響",
    "部署踩到的兩個坑:① EC2 role 少 s3:PutObject(500)② CI build amd64 但 EC2 arm64(502)→ buildx 跨平台修正",
], size=14, gap=9)

# ============ Slide 13: Status ============
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background(); bg.shadow.inherit = False
textbox(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(1.2), [
    ("原型狀態", 16, ACCENT, True, 4), ("可跑的 FastAPI 實作 + 管理頁前端", 30, WHITE, True)])
bullets_dark = [
    "7 個 API 端點 + 管理頁前端（建立/清單/編輯/刪除/分析）",
    "全部 curl 驗證通過：302 / 404 / 410 / 8 碼 token / 正規化 / SSRF / 過期",
    "production 替換點已在程式註解標明（Redis / 佇列 / CDN / 預聚合 / PostgreSQL）",
    "設計依據：DESIGN.md（16 題決策 + 優劣分析）",
]
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(3))
tf = tb.text_frame; tf.word_wrap = True
for i, it in enumerate(bullets_dark):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    r = p.add_run(); r.text = "✓  " + it
    _font(r, 16, RGBColor(0xE2, 0xE8, 0xF0), False)

# ============ Slide: 下週展開（Prototype → Production gap）============
s = prs.slides.add_slide(BLANK)
header(s, "下週預告 · NEXT WEEK", "從 Prototype 到 Production 的巨大鴻溝")
textbox(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
        [("快速掃過，建立全貌 — 下週 Deep Dive 才是主菜", 14, MUTED, False)])

gap_cards = [
    ("[Error Handling]", "遇到不合法輸入直接 Crash。", ANALYTICS),
    ("[Rate Limiting]", "毫無防護，極易被 Script 惡意灌爆 API。", ANALYTICS),
    ("[Auth & Isolation]", "缺乏多租戶概念，資料全公開。", ANALYTICS),
    ("[Monitoring]", "服務掛了完全沒有 Alerting，處於盲飛狀態。", GATEWAY),
    ("[Data Cleanup]", "過期資料永不刪除，一年後 DB 將面臨嚴重 Bloat。", GATEWAY),
    ("[Caching / CDN]", "每次轉址都直擊 DB，流量一來直接癱瘓。\n→ 下週 Deep Dive 展開", GATEWAY),
]
cw = Inches(3.95); ch = Inches(1.75); gx = Inches(0.27); gy = Inches(0.3)
x0 = Inches(0.6); y0 = Inches(2.15)
for i, (title, desc, accent) in enumerate(gap_cards):
    col = i % 3
    row = i // 3
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = LIGHT; card.shadow.inherit = False
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.09), ch)
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s, x + Inches(0.3), y + Inches(0.2), cw - Inches(0.5), ch - Inches(0.3),
            [(title, 16, INK, True, 8), (desc, 13, MUTED, False)])

# ============ Slide: 極限擴展 Step 1 — 瓶頸在哪 ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · STEP 1", "瓶頸在哪？—— DB 先爆還是 Server 先爆？")
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
cards = [
    ("API Server", ["Redirect handler 邏輯簡單（查 token → 回 302）", "CPU 不是瓶頸，水平擴展容易"], AMBER),
    ("Database", ["每次 redirect 都是一次 DB read", "50K QPS read → 大多數單機 DB 扛不住"], RED),
    ("Network", ["302 response 很小（只有 Location header）", "頻寬不是問題"], INK),
]
cy = Inches(1.75)
for title, lines, strip in cards:
    h = Inches(1.45)
    st = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), cy, Inches(0.12), h)
    st.fill.solid(); st.fill.fore_color.rgb = strip; st.line.fill.background(); st.shadow.inherit = False
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), cy, Inches(6.5), h)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT; bg.line.fill.background(); bg.shadow.inherit = False
    textbox(s, Inches(0.95), cy + Inches(0.12), Inches(6.1), h - Inches(0.2),
            [(title, 16, INK, True, 6)] + [(ln, 12.5, MUTED, False, 3) for ln in lines])
    cy += h + Inches(0.2)
concl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(2.1), Inches(5.15), Inches(2.9))
concl.fill.solid(); concl.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF); concl.line.color.rgb = ACCENT
concl.line.width = Pt(1.5); concl.shadow.inherit = False
textbox(s, Inches(7.9), Inches(2.4), Inches(4.6), Inches(2.4), [
    ("結論：DB 是第一個瓶頸", 18, ACCENT, True, 10),
    ("教材估算 ~5,800 QPS redirect", 13, MUTED, False, 4),
    ("現在要 50K → 約 10 倍", 13, INK, True, 12),
    ("解法方向：在 DB 前面擋一層 Cache", 14, INK, True)])
textbox(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.6), [
    ("讀側:cache 擋掉大多數 DB read（第 7 題 Redis 就是這層）→ 接著換「寫側(分析 scan_events)」變成下一個瓶頸(下頁)。", 12, MUTED, False)])

# ============ Slide: Step 2 — Cache hit rate math ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · STEP 2", "加 Cache 能解多少？—— 不同 hit rate 下 DB 承受量")
textbox(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5), [
    ("50K redirect QPS 全打 DB → 撐不住;cache 擋掉 hit 的部分,DB 只收 miss。", 13, INK, False)])
hrows = [
    ("Cache Hit Rate", "DB 實際 QPS", "能撐嗎？"),
    ("0%（沒 cache）", "50,000", "完全撐不住"),
    ("90%", "5,000", "勉強"),
    ("95%", "2,500", "可以"),
    ("99%", "500", "輕鬆"),
]
ht = s.shapes.add_table(len(hrows), 3, Inches(1.2), Inches(2.15), Inches(11.0), Inches(2.7)).table
ht.columns[0].width = Inches(4.0); ht.columns[1].width = Inches(3.5); ht.columns[2].width = Inches(3.5)
hcolor = [None, RED, ANALYTICS, ACCENT, GREEN]
for r in range(len(hrows)):
    for c in range(3):
        cell = ht.cell(r, c)
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4); cell.margin_left = Pt(10)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = hrows[r][c]
        col = WHITE if r == 0 else (hcolor[r] if c == 2 else INK)
        _font(run, 14, col, r == 0 or c == 2)
box(s, Inches(0.6), Inches(5.15), Inches(12.15), Inches(1.55),
    "活動 QR Code 特性:數量少（可能就幾十個）但每個 QPS 極高 → cache hit rate 可能 99%+ → DB 只收 ~500 QPS，完全可以。\n（行銷活動印在海報上的少數 QR，被幾百萬人掃 → 極端集中的熱點,天生適合 cache。）", DB, WHITE, 13)

# ============ Slide: Step 2 — Cache 選型 ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · STEP 2", "Cache 選型 — Local / Redis / CDN 三層比較")
crows2 = [
    ("維度", "Local Cache", "Redis / Memcached", "CDN Edge Cache"),
    ("延遲", "最低（in-process）", "中（network hop）", "最低（離用戶近）"),
    ("一致性", "差（每台各一份）", "好（共享）", "差（purge 慢）"),
    ("Hit Rate", "無 sticky routing 則低", "高（共享）", "依流量集中度"),
    ("容量", "受限於 server 記憶體", "獨立擴展", "依 CDN 供應商"),
    ("適用場景", "極熱 key、不常變", "通用", "靜態或少變內容"),
]
ct2 = s.shapes.add_table(len(crows2), 4, Inches(0.6), Inches(1.7), Inches(12.15), Inches(3.3)).table
ct2.columns[0].width = Inches(1.8); ct2.columns[1].width = Inches(3.45)
ct2.columns[2].width = Inches(3.45); ct2.columns[3].width = Inches(3.45)
for r in range(len(crows2)):
    for c in range(4):
        cell = ct2.cell(r, c)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = crows2[r][c]
        _font(run, 12, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
box(s, Inches(0.6), Inches(5.2), Inches(12.15), Inches(0.95),
    "本 scenario 選 Redis + CDN 雙層:CDN 擋第一層（離用戶近、減 origin 流量）· Redis 擋第二層（CDN miss 不直接打 DB）· DB 只處理 cold start / miss", GATEWAY, WHITE, 13)
textbox(s, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.7), [
    ("⚠ 雙層必要性看「規模 + 地理」:Local-only + 小流量 → 單層 Redis 也夠;跨區域 / >100K QPS → 雙層必要。本專案落點:第 7 題 Redis（單層）+ 附錄 B CDN 邊緣變體（選配升級到雙層）。", 12, MUTED, False, 2)])

# ============ Slide: Step 2b — Near Cache (Local + Redis) ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · STEP 2", "Near Cache（Local + Redis 兩層）何時用？")
textbox(s, Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.5), [
    ("核心動機:省掉「連 Redis 的那一跳」（同機房 ~0.5ms round-trip）。讀路徑加一層 L1:", 13, INK, False)])
for i, (t, c) in enumerate([("L1 local\n行程內 ~奈秒", CACHE), ("L2 Redis\n~0.5ms", GATEWAY), ("DB\ncache miss", DB)]):
    x = Inches(1.4) + i * Inches(3.6)
    box(s, x, Inches(2.05), Inches(2.9), Inches(0.9), t, c, WHITE, 12)
    if i < 2:
        arrow(s, x + Inches(2.9), Inches(2.5), x + Inches(3.6), Inches(2.5))
textbox(s, Inches(4.35), Inches(2.2), Inches(0.6), Inches(0.4), [("miss", 10, MUTED, False)])
textbox(s, Inches(7.95), Inches(2.2), Inches(0.6), Inches(0.4), [("miss", 10, MUTED, False)])
nrows = [
    ("值得加 L1 的情況", "為什麼要 L1"),
    ("極端熱點 key（hot key）", "少數 token 佔絕大多數流量（活動 QR 正是此型）→ L1 直接回、完全不碰 Redis"),
    ("單一 Redis 頻寬/CPU 快到頂", "50K+ QPS 打單節點會先飽和（附錄 J #4）→ L1 擋多數、替 Redis 卸流量"),
    ("要壓 p99 尾延遲", "Redis 偶發慢查詢/GC/網路抖動拉高 p99；L1 命中不受影響"),
]
nt = s.shapes.add_table(len(nrows), 2, Inches(0.6), Inches(3.2), Inches(12.15), Inches(2.0)).table
nt.columns[0].width = Inches(4.0); nt.columns[1].width = Inches(8.15)
for r in range(len(nrows)):
    for c in range(2):
        cell = nt.cell(r, c)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = nrows[r][c]
        _font(run, 12, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
box(s, Inches(0.6), Inches(5.45), Inches(12.15), Inches(0.85),
    "代價:L1 每台各一份 → 一致性變差（update 已在 Redis 失效,各台 L1 仍舊值）。解法:很短 L1 TTL（1–5s）或 pub/sub 廣播失效（keyspace notification）", RED, WHITE, 12)
textbox(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.6), [
    ("判斷:單層 Redis 已是瓶頸 + 流量集中 hot key + 能接受幾秒弱一致 → 才加 L1。兩種雙層:CDN+Redis（L1 放邊緣）vs Local+Redis（L1 放 app 行程內）。", 12, MUTED, False, 2)])

# ============ Slide: Step 3 — CDN TTL 兩難 ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · STEP 3", "CDN 的 TTL 兩難 + Analytics 缺口")
box(s, Inches(0.6), Inches(1.75), Inches(6.0), Inches(2.0),
    "TTL 太長（例 24h）\n\n• URL 改了，CDN 在 TTL 內不更新\n• 用戶還被導到舊目標頁\n• 無法即時停用被濫用的 short link", RED, WHITE, 13)
box(s, Inches(6.75), Inches(1.75), Inches(6.0), Inches(2.0),
    "TTL 太短（例 10s）\n\n• CDN 幾乎等於沒 cache\n• 大量流量穿透到 origin\n• 失去 CDN 擋流量的意義", ANALYTICS, WHITE, 13)
box(s, Inches(0.6), Inches(3.95), Inches(12.15), Inches(1.35),
    "解法:合理 TTL + 主動 invalidation\n主動 purge：URL 改時呼叫 CDN purge API（Cloudflare 秒級 / CloudFront 分鐘級）\nSWR (stale-while-revalidate)：max-age=300, swr=60 → 0–300s 回 cache；300–360s 回舊值+背景拉新（用戶 0 等待）；360s+ 同步 origin", ACCENT, WHITE, 12)
box(s, Inches(0.6), Inches(5.5), Inches(12.15), Inches(1.2),
    "⚠ CDN 的副作用:analytics 缺口 —— CDN 命中時 request 不打到 API Server → 這部分流量永遠不 emit analytics event → Analytics Pipeline 需額外手段補（CDN logs / edge function / 取樣）", RED, WHITE, 12)

# ============ Slide: Step 4 — Token 生成 scaling ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · STEP 4", "Token 生成也要 50K QPS？—— Creation flow 會被打爆嗎")
bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(3.0), [
    "Hash-Based（教材方案）",
    ("SHA-256 是 stateless 運算", 1),
    ("CPU bound → 水平擴展就行", 1),
    ("加機器 = 加產能，線性擴展", 1),
    ("Operational complexity 低", 1),
], size=14, gap=9)
bullets(s, Inches(6.9), Inches(1.9), Inches(6.0), Inches(3.0), [
    "Pre-generated Pool",
    ("事先產生一批 unique token 存 DB", 1),
    ("寫入時直接分配 → 零 collision", 1),
    ("Pool 會不會被耗盡？需背景 job 補充", 1),
    ("Operational complexity 高", 1),
], size=14, gap=9)
box(s, Inches(0.6), Inches(5.15), Inches(12.15), Inches(1.4),
    "判斷:活動場景 creation QPS 遠低於 redirect QPS（建一次、掃千萬次）→ Hash-based 足夠,不需 pre-generated pool 的額外複雜度。\n（token 生成不是這個 scenario 的瓶頸——瓶頸在 redirect 讀路徑,見 Step 1-3。）", DB, WHITE, 13)

# ============ Slide: 極限擴展 — 50x 推演 ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · 50× 推演", "1k → 50k QPS：這套架構哪裡先崩？")
textbox(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.6), [
    ("1,000 QPS  →  50,000 QPS（大型客戶上線,流量瞬間 50×）", 18, ACCENT, True)])
brows = [
    ("#", "元件", "為何崩（先 → 後）"),
    ("1", "DB 寫入 (scan_events)", "每次 redirect 一筆 INSERT → 50k writes/s 打單一 RDS（單 AZ、無 replica）→ 最先爆"),
    ("2", "EC2 ASG", "max_size=2 且無 scaling policy（固定 desired）→ 扛不住 50k 連線/CPU"),
    ("3", "API Gateway", "stage rate=1000 + 帳號預設 ~10k RPS → 大量 429"),
    ("4", "Redis 單節點", "cache.t4g.micro 50k GET/s 高 CPU/網路、無讀擴展"),
    ("5", "Edge", "CloudFront/ALB 可擴(LCU 成本↑);但 redirect 不快取 → 全額回源放大下游"),
]
bt = s.shapes.add_table(len(brows), 3, Inches(0.6), Inches(2.3), Inches(12.15), Inches(3.4)).table
bt.columns[0].width = Inches(0.7); bt.columns[1].width = Inches(3.3); bt.columns[2].width = Inches(8.15)
REDBG = RGBColor(0xFE, 0xE2, 0xE2)
for r in range(len(brows)):
    for c in range(3):
        cell = bt.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(8)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = ACCENT
        elif r == 1:
            cell.fill.fore_color.rgb = REDBG  # 最先崩
        else:
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = brows[r][c]
        _font(run, 12, WHITE if r == 0 else INK, r == 0 or c <= 1)
bullets(s, Inches(0.6), Inches(5.9), Inches(12.3), Inches(1.2), [
    "結論：DB 的 scan 寫入最先崩——分析寫入沒和 redirect 解耦(唯一線性暴增又集中單點的寫入)。",
], size=14, gap=4)

# ============ Slide: Event-Driven Analytics Pipeline ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · ANALYTICS PIPELINE", "解法:緩衝 + 批次聚合（解 DB 寫入瓶頸）")
box(s, Inches(0.5), Inches(2.1), Inches(2.1), Inches(0.9), "redirect\nput 1 event", CLIENT, WHITE, 12)
box(s, Inches(3.0), Inches(2.1), Inches(2.7), Inches(0.9), "Kinesis / SQS\n高吞吐緩衝(吸尖峰)", ANALYTICS, WHITE, 11)
box(s, Inches(6.1), Inches(2.1), Inches(2.7), Inches(0.9), "消費者\nLambda/KCL 批次聚合", APP, WHITE, 11)
box(s, Inches(9.2), Inches(1.4), Inches(3.3), Inches(0.7), "每日計數 rollup\n(DynamoDB/RDS)", DB, WHITE, 10)
box(s, Inches(9.2), Inches(2.25), Inches(3.3), Inches(0.7), "原始明細歸檔\n(S3 + Athena)", DB, WHITE, 10)
box(s, Inches(9.2), Inches(3.1), Inches(3.3), Inches(0.7), "即時分析(選)\nOpenSearch", DB, WHITE, 10)
arrow(s, Inches(2.6), Inches(2.55), Inches(3.0), Inches(2.55))
arrow(s, Inches(5.7), Inches(2.55), Inches(6.1), Inches(2.55))
arrow(s, Inches(8.8), Inches(2.55), Inches(9.2), Inches(1.75), color=MUTED, width=1.1)
arrow(s, Inches(8.8), Inches(2.55), Inches(9.2), Inches(2.6), color=MUTED, width=1.1)
arrow(s, Inches(8.8), Inches(2.55), Inches(9.2), Inches(3.45), color=MUTED, width=1.1)
bullets(s, Inches(0.6), Inches(4.4), Inches(12.3), Inches(2.6), [
    "redirect 只丟事件進 stream(極快、緩衝吸 50k/s 尖峰),不再每次同步寫 RDS → 寫入降到批次級(對應第 14 題)",
    "ASG：target-tracking autoscaling + 提高 max(或 Fargate/Lambda)",
    "讀路徑：Redis cluster/多節點 + 拉高熱門命中;RDS read replica 分攤 cache-miss 讀(第 16 題)",
    "edge：熱門 token 用 CloudFront 短 TTL 邊緣 redirect(附錄 B)在邊緣消化大半流量;API GW 調高 throttle",
], size=13, gap=8)

# ============ Slide: Analytics — 你想記錄什麼 ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · ANALYTICS", "你想記錄什麼？—— Who / When / Where / What")
awhat = [
    ("Who", "IP / User-Agent\n裝置類型", CLIENT),
    ("When", "掃描時間\n時區", GATEWAY),
    ("Where", "Referer / GeoIP\n來源頁面", CDN),
    ("What", "哪個 QR Code\n目標 URL", DB),
]
n = len(awhat); bw = Inches(2.9); gap = Inches(0.2); x0 = Inches(0.7); y = Inches(2.1)
for i, (t, d, color) in enumerate(awhat):
    box(s, x0 + i * (bw + gap), y, bw, Inches(1.5), t + "\n\n" + d, color, WHITE, 14)
box(s, Inches(0.7), Inches(4.15), Inches(12.0), Inches(1.05),
    "追問:這些資料若寫在 redirect 的 critical path 上 → 50K QPS 下每次多一個 DB write → DB 直接爆", RED, WHITE, 14)
bullets(s, Inches(0.7), Inches(5.4), Inches(12.0), Inches(1.4), [
    "所以拆兩條路徑:Redirect path（延遲敏感,要快 → Cache + DB）· Analytics path（量大可延遲,要穩 → Queue + Batch）",
    "本專案:scan_events 明細（第 8 題非同步寫）+ 正式版事件流（附錄 J pipeline）",
], size=13, gap=8)

# ============ Slide: Analytics Store 選型 ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · ANALYTICS", "Analytics Store 選型 —— 記錄下來存哪、怎麼查")
asrows = [
    ("", "OLTP (PostgreSQL)", "Time-Series (TimescaleDB)", "OLAP 自架 (ClickHouse)", "OLAP serverless (Athena/BigQuery)"),
    ("寫入", "簡單", "優化過", "極快（batch）", "便宜"),
    ("趨勢查詢", "慢", "快（時間分區）", "極快（columnar）", "中"),
    ("單筆查詢", "快（索引）", "中", "中", "慢"),
    ("成本", "高", "中", "中", "最低"),
    ("適合", "小規模", "中規模+即時", "大規模+即時", "大規模+報表"),
]
at = s.shapes.add_table(len(asrows), 5, Inches(0.5), Inches(1.75), Inches(12.35), Inches(3.2)).table
at.columns[0].width = Inches(1.5)
for ci in (1, 2, 3, 4):
    at.columns[ci].width = Inches(2.71)
for r in range(len(asrows)):
    for c in range(5):
        cell = at.cell(r, c)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2); cell.margin_left = Pt(6)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = asrows[r][c]
        _font(run, 11, WHITE if r == 0 else (INK if c == 0 else MUTED), r == 0 or c == 0)
box(s, Inches(0.5), Inches(5.2), Inches(12.35), Inches(1.35),
    "選型取決於查詢模式:即時 dashboard + 時間序列 → TimescaleDB｜即時 + 複雜聚合 → ClickHouse 自架\n每日/每週報表 + 不想維運 → Athena / BigQuery 最划算（本專案附錄 J:S3 + Athena 歸檔 + RDS/DynamoDB 預聚合）", DB, WHITE, 12)

# ============ Slide: 即時 vs 每日報表 ============
s = prs.slides.add_slide(BLANK)
header(s, "極限擴展 · ANALYTICS", "即時 dashboard vs 每日報表 —— 架構一樣嗎？")
textbox(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.4), [
    ("即時 Dashboard（streaming，延遲 < 秒級,複雜度/成本高）", 14, ACCENT, True)])
for i, (t, c) in enumerate([("Kafka", ANALYTICS), ("Flink / Spark\nStreaming", ANALYTICS), ("TimescaleDB", ANALYTICS), ("Dashboard\n(即時更新)", GATEWAY)]):
    box(s, Inches(0.7) + i * Inches(3.05), Inches(2.05), Inches(2.7), Inches(0.9), t, c, WHITE, 12)
    if i < 3:
        arrow(s, Inches(0.7) + i * Inches(3.05) + Inches(2.7), Inches(2.5), Inches(0.7) + (i + 1) * Inches(3.05), Inches(2.5))
textbox(s, Inches(0.6), Inches(3.3), Inches(12), Inches(0.4), [
    ("每日報表（batch，延遲 = 小時級,簡單便宜）", 14, GREEN, True)])
for i, (t, c) in enumerate([("SQS", GREEN), ("Consumer\n→ S3", GREEN), ("Daily\nBatch Job", GREEN), ("Report\n(每日更新)", GATEWAY)]):
    box(s, Inches(0.7) + i * Inches(3.05), Inches(3.8), Inches(2.7), Inches(0.9), t, c, WHITE, 12)
    if i < 3:
        arrow(s, Inches(0.7) + i * Inches(3.05) + Inches(2.7), Inches(4.25), Inches(0.7) + (i + 1) * Inches(3.05), Inches(4.25))
box(s, Inches(0.6), Inches(5.1), Inches(12.15), Inches(1.15),
    "面試重點:先問需求再選架構。大部分 QR Code 服務「每日報表就夠了」→ 不要過度設計。\n（需要秒級即時才上 Kafka+Flink;否則 SQS + S3 + 每日 batch 最省。）", ACCENT, WHITE, 13)

# ============ Slide: 面試框架 — 時間分配 ============
s = prs.slides.add_slide(BLANK)
header(s, "面試實戰 · FRAMEWORK", "Design a QR Code Generator —— 45 分鐘時間分配")
iframe = [
    ("0–5 min：需求釐清", "靜態 or 動態？需要 analytics 嗎？預期 QPS？→ 問對問題 = 第一個加分點", AMBER),
    ("5–15 min：High-Level Design", "教材內容要能流暢畫完:Creation flow + Retrieval flow", ACCENT),
    ("15–35 min：Deep Dive（面試官挑 1–2 追問）", "Token collision → Week 1｜流量暴增 → Block 2｜Analytics → Block 3", DB),
    ("35–45 min：收尾", "總結 trade-off + 反問面試官問題", GREEN),
]
cy = Inches(1.85)
for title, body, strip in iframe:
    st = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), cy, Inches(0.12), Inches(1.0))
    st.fill.solid(); st.fill.fore_color.rgb = strip; st.line.fill.background(); st.shadow.inherit = False
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), cy, Inches(12.0), Inches(1.0))
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT; bg.line.fill.background(); bg.shadow.inherit = False
    textbox(s, Inches(0.95), cy + Inches(0.12), Inches(11.6), Inches(0.8),
            [(title, 15, INK, True, 5), (body, 12.5, MUTED, False)])
    cy += Inches(1.2)

# ============ Slide: 面試常見扣分點 ============
s = prs.slides.add_slide(BLANK)
header(s, "面試實戰 · PITFALLS", "面試常見扣分點")
pitfalls = [
    ("需求沒問就開始畫圖", "靜態 vs 動態都沒確認,直接假設 → 面試官質疑溝通能力"),
    ("只講 Happy Path", "沒考慮 failure case:cache miss、DB 掛、token collision → 缺工程深度"),
    ("Deep Dive 太淺", "每個都講一點但都不深 → 不如挑一個講透,展示能「往下鑽」"),
    ("沒有量化分析", "只說「加 cache 就好」但不算 hit rate、不算 DB 負載 → 缺說服力"),
]
cy = Inches(1.9)
for title, body in pitfalls:
    st = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), cy, Inches(0.12), Inches(1.0))
    st.fill.solid(); st.fill.fore_color.rgb = RED; st.line.fill.background(); st.shadow.inherit = False
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), cy, Inches(12.0), Inches(1.0))
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT; bg.line.fill.background(); bg.shadow.inherit = False
    textbox(s, Inches(0.95), cy + Inches(0.14), Inches(11.6), Inches(0.8),
            [(title, 15, INK, True, 5), (body, 12.5, MUTED, False)])
    cy += Inches(1.25)

# ============ Slide: 現場練習 — cache hit 掉到 50% ============
s = prs.slides.add_slide(BLANK)
header(s, "面試實戰 · PRACTICE", "追問:Cache hit rate 掉到 50%，你怎麼辦？")
prrows = [
    ("診斷方向", "應對策略"),
    ("Working set > cache size（QR 數量暴增）", "加大 Redis 記憶體 / Redis Cluster 分片"),
    ("TTL 設太短 → 頻繁過期 miss", "延長 TTL / stale-while-revalidate"),
    ("Eviction policy 不對（LRU vs LFU）", "有 hot key → 換 LFU（Redis 4.0+）"),
    ("Bot 掃冷門 token → cache penetration", "Negative caching：404 也 cache（短 TTL）"),
]
pr = s.shapes.add_table(len(prrows), 2, Inches(0.6), Inches(2.0), Inches(12.15), Inches(3.0)).table
pr.columns[0].width = Inches(5.9); pr.columns[1].width = Inches(6.25)
for r in range(len(prrows)):
    for c in range(2):
        cell = pr.cell(r, c)
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4); cell.margin_left = Pt(10)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (WHITE if r % 2 else LIGHT)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = prrows[r][c]
        _font(run, 13, WHITE if r == 0 else INK, r == 0)
box(s, Inches(0.6), Inches(5.35), Inches(12.15), Inches(1.0),
    "答題要點:先「診斷」再「對策」—— 分辨是容量、TTL、eviction、還是穿透問題,對症下藥,而非一律「加機器」。", DB, WHITE, 13)

prs.save("QR_Code_Generator.pptx")
print("saved QR_Code_Generator.pptx ·", len(prs.slides._sldIdLst), "slides")
